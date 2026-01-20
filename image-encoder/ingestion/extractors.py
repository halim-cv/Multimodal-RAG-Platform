"""
Image extraction from various sources
"""
import shutil
import zipfile
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
import mimetypes

try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    logging.warning("PyMuPDF not available - PDF extraction disabled")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logging.warning("python-pptx not available - PPTX extraction disabled")

try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False
    logging.warning("requests not available - URL download disabled")

import config
from preprocessing import validate_image

logger = logging.getLogger(__name__)


class ExtractedImageRecord:
    """Data class for extracted images"""
    def __init__(
        self,
        source: str,
        source_type: str,
        temp_path: Path,
        source_index: Optional[int] = None
    ):
        self.source = source
        self.source_type = source_type
        self.temp_path = temp_path
        self.source_index = source_index
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "source": self.source,
            "source_type": self.source_type,
            "temp_path": str(self.temp_path),
            "source_index": self.source_index
        }


def extract_from_file(src_path: Path) -> List[ExtractedImageRecord]:
    """
    Extract images from a file (PDF, PPTX, image, ZIP)
    
    Args:
        src_path: Path to source file
    
    Returns:
        List of ExtractedImageRecord
    """
    if not src_path.exists():
        logger.error(f"File not found: {src_path}")
        return []
    
    # Detect file type
    mime_type, _ = mimetypes.guess_type(str(src_path))
    suffix = src_path.suffix.lower()
    
    logger.info(f"Extracting from {src_path.name} (type: {mime_type or suffix})")
    
    # Route to appropriate extractor
    if suffix == ".pdf":
        return extract_from_pdf(src_path)
    elif suffix in [".pptx", ".ppt"]:
        return extract_from_pptx(src_path)
    elif suffix == ".zip":
        return extract_from_zip(src_path)
    elif suffix in [".jpg", ".jpeg", ".png", ".webp", ".tiff", ".bmp", ".gif"]:
        return extract_raw_image(src_path)
    else:
        logger.warning(f"Unsupported file type: {suffix}")
        return []


def extract_from_pdf(pdf_path: Path) -> List[ExtractedImageRecord]:
    """Extract images from PDF using PyMuPDF"""
    if not HAS_PYMUPDF:
        logger.error("PyMuPDF not installed - cannot extract from PDF")
        return []
    
    logger.info(f"Extracting from PDF: {pdf_path.name}")
    extracted = []
    
    try:
        doc = fitz.open(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Extract embedded images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                # Save to temp file
                filename = f"{pdf_path.stem}_page{page_num+1}_img{img_index+1}.{image_ext}"
                temp_path = config.RAW_DIR / filename
                
                with open(temp_path, "wb") as f:
                    f.write(image_bytes)
                
                logger.debug(f"Extracted embedded image: {filename}")
                
                if validate_image(temp_path):
                    extracted.append(ExtractedImageRecord(
                        source=str(pdf_path),
                        source_type="pdf",
                        temp_path=temp_path,
                        source_index=page_num
                    ))
                else:
                    temp_path.unlink()
            
            # Also render page as image (alternative method)
            # Uncomment if you want full page renders instead
            # pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
            # filename = f"{pdf_path.stem}_page{page_num+1}_render.png"
            # temp_path = config.RAW_DIR / filename
            # pix.save(str(temp_path))
            # extracted.append(ExtractedImageRecord(...))
        
        doc.close()
        logger.info(f"Extracted {len(extracted)} images from PDF")
        
    except Exception as e:
        logger.error(f"Failed to extract from PDF {pdf_path.name}: {e}")
    
    return extracted


def extract_from_pptx(pptx_path: Path) -> List[ExtractedImageRecord]:
    """Extract images from PowerPoint"""
    if not HAS_PPTX:
        logger.error("python-pptx not installed - cannot extract from PPTX")
        return []
    
    logger.info(f"Extracting from PPTX: {pptx_path.name}")
    extracted = []
    
    try:
        prs = Presentation(pptx_path)
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if shape is a picture
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    
                    # Get image bytes
                    image_bytes = image.blob
                    image_ext = image.ext
                    
                    # Save to temp file
                    filename = f"{pptx_path.stem}_slide{slide_idx+1}_img{shape_idx+1}.{image_ext}"
                    temp_path = config.RAW_DIR / filename
                    
                    with open(temp_path, "wb") as f:
                        f.write(image_bytes)
                    
                    logger.debug(f"Extracted PPTX image: {filename}")
                    
                    if validate_image(temp_path):
                        extracted.append(ExtractedImageRecord(
                            source=str(pptx_path),
                            source_type="pptx",
                            temp_path=temp_path,
                            source_index=slide_idx
                        ))
                    else:
                        temp_path.unlink()
        
        logger.info(f"Extracted {len(extracted)} images from PPTX")
        
    except Exception as e:
        logger.error(f"Failed to extract from PPTX {pptx_path.name}: {e}")
    
    return extracted


def extract_from_zip(zip_path: Path) -> List[ExtractedImageRecord]:
    """Extract images from ZIP file"""
    logger.info(f"Extracting from ZIP: {zip_path.name}")
    extracted = []
    
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            for file_info in zip_ref.filelist:
                filename = file_info.filename
                
                # Check if it's an image file
                if any(filename.lower().endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.webp', '.tiff', '.bmp']):
                    # Extract to temp location
                    temp_path = config.RAW_DIR / Path(filename).name
                    
                    with zip_ref.open(file_info) as source, open(temp_path, 'wb') as target:
                        shutil.copyfileobj(source, target)
                    
                    logger.debug(f"Extracted from ZIP: {filename}")
                    
                    if validate_image(temp_path):
                        extracted.append(ExtractedImageRecord(
                            source=str(zip_path),
                            source_type="zip",
                            temp_path=temp_path,
                            source_index=None
                        ))
                    else:
                        temp_path.unlink()
        
        logger.info(f"Extracted {len(extracted)} images from ZIP")
        
    except Exception as e:
        logger.error(f"Failed to extract from ZIP {zip_path.name}: {e}")
    
    return extracted


def extract_raw_image(image_path: Path) -> List[ExtractedImageRecord]:
    """Copy raw image to raw directory"""
    logger.info(f"Processing raw image: {image_path.name}")
    
    if not validate_image(image_path):
        logger.warning(f"Invalid image: {image_path.name}")
        return []
    
    # Copy to raw directory if not already there
    temp_path = config.RAW_DIR / image_path.name
    
    # Check if source and dest are the same file to avoid SameFileError
    try:
        if image_path.resolve() != temp_path.resolve():
            shutil.copy2(image_path, temp_path)
            logger.debug(f"Copied to raw: {temp_path.name}")
        else:
            logger.debug(f"File already in raw dir: {temp_path.name}")
    except OSError:
        # Fallback for systems where resolve() might behave unexpectedly
        if image_path.absolute() != temp_path.absolute():
            shutil.copy2(image_path, temp_path)

    
    return [ExtractedImageRecord(
        source=str(image_path),
        source_type="image",
        temp_path=temp_path,
        source_index=None
    )]


def download_from_url(url: str, output_dir: Path = None) -> Optional[Path]:
    """
    Download image from URL
    
    Args:
        url: Image URL
        output_dir: Directory to save (defaults to RAW_DIR)
    
    Returns:
        Path to downloaded file or None
    """
    if not HAS_REQUESTS:
        logger.error("requests not installed - cannot download from URL")
        return None
    
    if output_dir is None:
        output_dir = config.RAW_DIR
    
    logger.info(f"Downloading from URL: {url}")
    
    try:
        response = requests.get(url, timeout=30, stream=True)
        response.raise_for_status()
        
        # Determine filename
        filename = Path(url).name
        if not filename or '.' not in filename:
            filename = f"download_{hash(url) % 100000}.jpg"
        
        output_path = output_dir / filename
        
        with open(output_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        logger.info(f"Downloaded to: {output_path.name}")
        return output_path
        
    except Exception as e:
        logger.error(f"Failed to download from {url}: {e}")
        return None
